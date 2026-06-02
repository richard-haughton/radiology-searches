#!/usr/bin/env python3
"""
Extract findings from 'Diagnostic Radiology On Call Bootcamp.pptx'
and generate scripts/data/bootcamp_radiology_findings_seed.json.

Usage:
  conda activate searches && python scripts/extract_bootcamp_findings.py
"""

import base64
import io
import json
import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

PPTX_PATH = os.path.join(
    os.path.dirname(__file__), '..', 'Diagnostic Radiology On Call Bootcamp.pptx'
)
OUT_PATH = os.path.join(
    os.path.dirname(__file__), 'data', 'bootcamp_radiology_findings_seed.json'
)

# ---------------------------------------------------------------------------
# Slide-number → finding mapping
#
# Each entry defines ONE finding.  'slides' lists the 1-based slide numbers
# whose text and images feed this finding.  Multiple consecutive slides that
# belong to the same topic are merged.
#
# patternName / stepTitle tell the import script where to link the finding
# inside the user's search-pattern library.
# ---------------------------------------------------------------------------

FINDINGS_SPEC = [
    # ── NEURO / CT HEAD ──────────────────────────────────────────────────
    {
        "name": "Subarachnoid Hemorrhage",
        "slides": [10, 11],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "6. Examine the CSF spaces.",
        "extra": {
            "clinical": "Sudden severe 'thunderclap' headache; photophobia, neck stiffness. May be traumatic or aneurysmal.",
            "report": "Hyperdensity in the sulci/basal cisterns consistent with subarachnoid hemorrhage. CTA recommended to exclude aneurysmal source.",
            "treatment": "ICU admission; nimodipine; neurosurgery/interventional radiology consult.",
        },
    },
    {
        "name": "Subdural Hemorrhage",
        "slides": [12],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "6. Examine the CSF spaces.",
        "extra": {
            "clinical": "Head trauma; altered mental status. Crosses suture lines; crescentic shape.",
            "report": "Crescent-shaped extra-axial collection along the (right/left) cerebral convexity consistent with (acute/subacute/chronic) subdural hematoma. (Mass effect/midline shift: X mm.)",
            "treatment": "Neurosurgery consult; may require burr-hole drainage or craniotomy depending on size and symptoms.",
        },
    },
    {
        "name": "Epidural Hematoma",
        "slides": [13],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "6. Examine the CSF spaces.",
        "extra": {
            "clinical": "Classic 'talk and die' — brief LOC, lucid interval, rapid deterioration. Does NOT cross suture lines; biconvex/lens-shaped.",
            "report": "Biconvex hyperdense extra-axial collection adjacent to the (temporal/parietal/frontal) bone consistent with acute epidural hematoma.",
            "treatment": "Emergent neurosurgery consult; craniotomy for hematoma evacuation.",
        },
    },
    {
        "name": "Intraparenchymal Hematoma",
        "slides": [14],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "8. Assess the brain parenchyma for intra-parenchymal blood/ mass lesion.",
        "extra": {
            "clinical": "Acute neurologic deficit with or without trauma; hypertensive hemorrhage classically in basal ganglia/thalamus.",
            "report": "Intraparenchymal hyperdensity in the (location) consistent with acute hemorrhage. Note: globus pallidus mineralization is a mimic — no vasogenic edema.",
            "treatment": "ICU admission; aggressive BP control; coagulopathy reversal; neurosurgery consult for significant mass effect.",
        },
    },
    {
        "name": "MRI Blood Product Aging (I Bleed Die)",
        "slides": [15],
        "modalities": ["MRI"],
        "isRedFinding": False,
        "patternName": "MRI Brain",
        "stepTitle": "6. Examine the SWI sequence.",
        "extra": {
            "clinical": "Aide-mémoire for MRI signal of intracranial blood: Isointense (hyperacute) → Bright (acute–subacute) → Dark (chronic).",
            "report": "MRI signal characteristics of blood products reflect age: hyperacute isointense; acute T1 iso/T2 dark; early subacute T1 bright; late subacute T1/T2 bright; chronic T1/T2 dark (hemosiderin).",
            "treatment": "N/A — reference tool for interpreting MRI hemorrhage age.",
        },
    },
    {
        "name": "Cerebral Herniation",
        "slides": [17],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "9. Look for mass effect/herniation.",
        "extra": {
            "clinical": "Depressed consciousness, anisocoria, Cushing's triad. Complication of mass lesion, edema, or hemorrhage.",
            "report": "Subfalcine herniation with X mm midline shift. Transalar (ascending/descending) herniation. Uncal herniation with effacement of the ipsilateral suprasellar cistern.",
            "treatment": "Emergent neurosurgery; mannitol/hypertonic saline; intubation; ICP monitoring.",
        },
    },
    {
        "name": "Skull Fractures",
        "slides": [18],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "5. Examine the bones, look for fractures, aggressive lesions.",
        "extra": {
            "clinical": "Head trauma. Describe: comminuted vs. simple; suture involvement; depressed vs. non-depressed; temporal bone involvement (CN VII risk).",
            "report": "Fracture through the (location); (depressed/non-depressed); (crosses/does not cross) suture lines; (involves/spares) the temporal bone.",
            "treatment": "Neurosurgery for depressed fractures over eloquent cortex; ENT if temporal bone involved; observation for simple non-depressed.",
        },
    },
    {
        "name": "Carotid/Vertebral Vascular Injury (Denver/Biffl Scale)",
        "slides": [19, 20, 21, 22],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "Look at the vasculature",
        "extra": {
            "clinical": "High-energy neck trauma. Denver/Biffl scale grades I–V; pseudo-occlusion can mimic complete occlusion.",
            "report": "Grade (I–V) injury of the (right/left) internal carotid/vertebral artery: intimal irregularity / pseudoaneurysm / complete occlusion. Correlation with CTA/DSA recommended.",
            "treatment": "Antithrombotic therapy; endovascular or surgical repair for higher-grade injuries.",
        },
    },
    {
        "name": "Floating Thrombus",
        "slides": [23],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "Look at the vasculature",
        "extra": {
            "clinical": "Freely mobile thrombus within a vessel lumen; high stroke risk.",
            "report": "Filling defect within the (vessel) consistent with a floating/pedunculated thrombus.",
            "treatment": "Anticoagulation; vascular surgery or interventional radiology consult.",
        },
    },
    {
        "name": "Carotid Webs",
        "slides": [24],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CT Head",
        "stepTitle": "Look at the vasculature",
        "extra": {
            "clinical": "Young patient with cryptogenic stroke; thin shelf-like filling defect in the carotid bulb posteriorly.",
            "report": "Thin shelf-like filling defect at the posterior wall of the carotid bulb consistent with a carotid web.",
            "treatment": "Antiplatelet therapy; consider surgical endarterectomy or stenting.",
        },
    },
    {
        "name": "Fibromuscular Dysplasia (FMD)",
        "slides": [25],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CT Head",
        "stepTitle": "Look at the vasculature",
        "extra": {
            "clinical": "Young female; absence of other traumatic injuries. Classic 'string of beads' or smooth narrowing pattern.",
            "report": "Alternating stenoses and dilatations of the (ICA/vertebral artery) with 'string of beads' appearance consistent with fibromuscular dysplasia.",
            "treatment": "Antiplatelet therapy; angioplasty for symptomatic cases; long-term surveillance.",
        },
    },
    {
        "name": "Ischemic Stroke – Arterial Territory",
        "slides": [26],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "10. Look specifically for signs of ischemic stroke.",
        "extra": {
            "clinical": "Acute focal neurologic deficit in an arterial territory (FAST). Some trauma patients had a stroke that caused the event.",
            "report": "Subtle hypoattenuation in the (MCA/ACA/PCA) territory with (loss of grey-white differentiation/hyperdense vessel sign) consistent with acute ischemic stroke. CTA for vessel occlusion.",
            "treatment": "Neurology consult; thrombolysis eligibility assessment; mechanical thrombectomy for large vessel occlusion.",
        },
    },
    {
        "name": "Ischemic Stroke – Venous (Cerebral Venous Thrombosis)",
        "slides": [27],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "10. Look specifically for signs of ischemic stroke.",
        "extra": {
            "clinical": "Non-arterial territory; often hemorrhagic; younger patients. Suspect in peri-partum, hypercoagulable states, OCP use.",
            "report": "Hemorrhagic infarct in non-arterial distribution; (hyperdense sinus sign) consistent with cerebral venous thrombosis. CT venography recommended.",
            "treatment": "Anticoagulation even in hemorrhagic transformation; hematology consult.",
        },
    },
    {
        "name": "Intracranial Aneurysms",
        "slides": [28],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "Look at the vasculature",
        "extra": {
            "clinical": "Incidental or presenting with SAH. Common locations: ACOM, PCOM, MCA bifurcation, basilar tip.",
            "report": "Saccular aneurysm at the (location) measuring X mm. Report neck, dome, relationship to parent vessel, and any perianeurysmal blood.",
            "treatment": "Neurosurgery/interventional neuroradiology consult; clipping vs. coiling based on morphology.",
        },
    },
    {
        "name": "Intracranial Infection (Cerebritis / Empyema / Abscess / Meningitis / Ventriculitis)",
        "slides": [29],
        "modalities": ["MRI", "CT"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "8. Assess the brain parenchyma for intra-parenchymal blood/ mass lesion.",
        "extra": {
            "clinical": "Fever, headache, altered mental status; recent ear/sinus infection or neurosurgery.",
            "report": "Ring-enhancing lesion with restricted diffusion in the (location) consistent with pyogenic abscess. Adjacent leptomeningeal enhancement suggests meningitis/ventriculitis component.",
            "treatment": "IV antibiotics; infectious disease and neurosurgery consult; consider aspiration/drainage.",
        },
    },
    {
        "name": "Chronic Microvascular Ischemia and Volume Loss",
        "slides": [30, 31],
        "modalities": ["CT", "MRI"],
        "isRedFinding": False,
        "patternName": "CT Head",
        "stepTitle": "8. Assess the brain parenchyma for intra-parenchymal blood/ mass lesion.",
        "extra": {
            "clinical": "Age-related finding in elderly with vascular risk factors. Beware of overcalling lacunar infarcts by age.",
            "report": "Scattered periventricular/subcortical white matter hypodensities consistent with chronic small vessel ischemic disease. Encephalomalacia (volume loss with CSF replacement) in (location).",
            "treatment": "Risk factor modification; no acute intervention typically required.",
        },
    },
    {
        "name": "Intracranial Malignancy (Primary and Metastatic)",
        "slides": [32],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "8. Assess the brain parenchyma for intra-parenchymal blood/ mass lesion.",
        "extra": {
            "clinical": "Headache, focal deficits, seizures; known primary malignancy raises suspicion for mets.",
            "report": "Ring-enhancing lesion(s) with vasogenic edema in the (location). Differential includes metastasis vs. high-grade glioma vs. abscess. MRI with contrast recommended.",
            "treatment": "Neurosurgery, oncology, radiation oncology consult; steroids for edema management.",
        },
    },
    {
        "name": "Catastrophic Brain Injury",
        "slides": [33],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "8. Assess the brain parenchyma for intra-parenchymal blood/ mass lesion.",
        "extra": {
            "clinical": "Non-survivable or near non-survivable injury. Diffuse oedema, herniation, absent basal cisterns. 'Sometimes these cases are so complicated they are easy.'",
            "report": "Diffuse cerebral edema with effacement of sulci/cisterns and (transtentorial/cerebellar tonsillar) herniation consistent with catastrophic brain injury.",
            "treatment": "ICU/palliative consult; family meeting; organ donation consideration.",
        },
    },
    {
        "name": "Demyelination",
        "slides": [34],
        "modalities": ["CT", "MRI"],
        "isRedFinding": False,
        "patternName": "CT Head",
        "stepTitle": "8. Assess the brain parenchyma for intra-parenchymal blood/ mass lesion.",
        "extra": {
            "clinical": "Young patient with relapsing-remitting neurologic symptoms; white matter hypoattenuating areas on CT.",
            "report": "Multiple periventricular/juxtacortical T2/FLAIR hyperintense lesions in a distribution typical of demyelinating disease (MS). Ovoid lesions perpendicular to the ventricular surface ('Dawson fingers').",
            "treatment": "Neurology consult; disease-modifying therapy if MS confirmed.",
        },
    },
    {
        "name": "Anoxic Brain Injury",
        "slides": [35],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "8. Assess the brain parenchyma for intra-parenchymal blood/ mass lesion.",
        "extra": {
            "clinical": "Post-cardiac arrest or severe hypoxia. Symmetric cortical and deep grey nucleus involvement.",
            "report": "Diffuse symmetric cortical hypoattenuation with loss of grey-white differentiation and deep grey nucleus involvement consistent with hypoxic-ischemic encephalopathy. MRI: symmetric DWI/ADC abnormality throughout cortex and deep structures.",
            "treatment": "Targeted temperature management; neurology/critical care consult; prognostication.",
        },
    },
    {
        "name": "Fat Embolism Syndrome (Brain)",
        "slides": [36],
        "modalities": ["MRI", "CT"],
        "isRedFinding": True,
        "patternName": "CT Head",
        "stepTitle": "8. Assess the brain parenchyma for intra-parenchymal blood/ mass lesion.",
        "extra": {
            "clinical": "Long bone fracture; classic triad of respiratory, neurologic, and petechial findings 24-72 h after injury.",
            "report": "Diffuse punctate DWI-restricted lesions in a 'starfield' pattern involving white matter and corpus callosum consistent with fat embolism syndrome.",
            "treatment": "Supportive; corticosteroids controversial; orthopedic stabilization.",
        },
    },

    # ── SPINE ─────────────────────────────────────────────────────────────
    {
        "name": "Spine Trauma – Occipito-Atlantal and C1-C2 Instability",
        "slides": [37, 38],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Cervical Spine",
        "stepTitle": "9. Assess the bones on sagittal images.",
        "extra": {
            "clinical": "High-energy mechanism; OA and C1-C2 injuries frequently missed on call. Evaluate all three planes.",
            "report": "Evaluate atlanto-occipital and atlanto-axial alignment. Describe displacement and canal compromise at C1-C2.",
            "treatment": "Spine surgery consult; halo immobilization vs. surgical stabilization.",
        },
    },
    {
        "name": "Burst Fracture",
        "slides": [39],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Spine T/L",
        "stepTitle": "9. Assess the osseous vertebra and discs.",
        "extra": {
            "clinical": "Axial loading injury. Violates dorsal cortex; retropulsion narrows spinal canal.",
            "report": "Burst fracture at (level) with retropulsion of posterior cortex, narrowing the spinal canal by approximately X%. MRI to evaluate cord signal if neurologically symptomatic.",
            "treatment": "Spine surgery consult; orthosis vs. surgical stabilization based on degree of instability and neurologic status.",
        },
    },
    {
        "name": "GRE/SWI Blooming for Spinal Hemorrhage",
        "slides": [40],
        "modalities": ["MRI"],
        "isRedFinding": True,
        "patternName": "MRI Thoracic Spine",
        "stepTitle": "4. Assess the sagittal and axial T2 images.",
        "extra": {
            "clinical": "Spinal cord hemorrhage or hematomyelia; GRE/SWI sequences 'bloom' (exaggerate) blood products.",
            "report": "Blooming susceptibility artifact on GRE/SWI at (level) consistent with intramedullary/extramedullary hemorrhage.",
            "treatment": "Spine surgery consult; MRI follow-up for cord injury assessment.",
        },
    },
    {
        "name": "Chance Fracture",
        "slides": [41],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Spine T/L",
        "stepTitle": "9. Assess the osseous vertebra and discs.",
        "extra": {
            "clinical": "Flexion-distraction injury (lap-belt); compression of anterior column + distraction of middle/posterior columns. High association with intra-abdominal injury (duodenal laceration).",
            "report": "Chance fracture at (level) with anterior compression and posterior distraction component. Correlate with abdomen for associated injury.",
            "treatment": "Thoracolumbar brace vs. surgical fixation; screen for abdominal injuries.",
        },
    },
    {
        "name": "Facet Injury (Unilateral/Bilateral Locked Facets)",
        "slides": [42],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Spine T/L",
        "stepTitle": "10. Assess the bones on axial images.",
        "extra": {
            "clinical": "Flexion-distraction mechanism; bilateral locked facets have 100% cord injury risk.",
            "report": "Unilateral/bilateral facet dislocation at (level) with perched/locked facets. Evaluate canal compromise and neural foraminal narrowing.",
            "treatment": "Spine surgery consult; traction reduction prior to fixation.",
        },
    },
    {
        "name": "Fused Spine (DISH / Ankylosing Spondylitis) Fractures",
        "slides": [43],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Spine T/L",
        "stepTitle": "9. Assess the osseous vertebra and discs.",
        "extra": {
            "clinical": "Brittle fused spine; fractures often missed; propagate through disc spaces ('three-column injury') even with minor trauma.",
            "report": "Horizontal fracture through a fused/bridging segment at (level) in the setting of (DISH/ankylosing spondylitis). Unstable injury.",
            "treatment": "Urgent spine surgery consult; prone to epidural hematoma formation; MRI recommended.",
        },
    },
    {
        "name": "Acute Disc Herniation",
        "slides": [44],
        "modalities": ["CT", "MRI"],
        "isRedFinding": False,
        "patternName": "CT Spine T/L",
        "stepTitle": "9. Assess the osseous vertebra and discs.",
        "extra": {
            "clinical": "Not all disc herniations are chronic; acute herniation in context of trauma may require urgent intervention.",
            "report": "Disc herniation at (level) with (central/paracentral/foraminal) component causing (X mm) canal/foraminal narrowing. Relevant in context of acute neurologic symptoms.",
            "treatment": "Pain management; neurosurgery consult for acute cord compression or cauda equina syndrome.",
        },
    },
    {
        "name": "Sacral Fractures",
        "slides": [45],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Spine T/L",
        "stepTitle": "9. Assess the osseous vertebra and discs.",
        "extra": {
            "clinical": "High-energy pelvic trauma; frequently missed. H-type sacral fracture has high neurologic risk.",
            "report": "Sacral fracture at (zone I/II/III) with (involvement/sparing) of sacral foramina and neural canal.",
            "treatment": "Trauma surgery and spine surgery consult; percutaneous sacroiliac screw fixation.",
        },
    },
    {
        "name": "Spinal Epidural Hematoma",
        "slides": [46],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Spine T/L",
        "stepTitle": "5. Do a quick first look at the spinal canal.",
        "extra": {
            "clinical": "Acute back pain with progressive myelopathy; spontaneous (anticoagulation) or post-traumatic.",
            "report": "Epidural hyperdensity (CT) / T1-bright/T2-variable collection (MRI) at (level) with cord compression and myelopathic signal change.",
            "treatment": "Emergent spine surgery consult for decompression; reverse anticoagulation.",
        },
    },
    {
        "name": "Spine Infection (Osteomyelitis / Discitis / Epidural Abscess)",
        "slides": [47],
        "modalities": ["MRI", "CT"],
        "isRedFinding": True,
        "patternName": "CT Spine T/L",
        "stepTitle": "14. Look at the paraspinal soft tissues.",
        "extra": {
            "clinical": "Back pain with fever; IV drug use, recent spine procedure, or immunocompromised state.",
            "report": "Disc and endplate destruction at (level) with paravertebral soft tissue thickening/fluid and epidural enhancement consistent with discitis-osteomyelitis. Epidural phlegmon/abscess (if present) at (level) causing cord compression.",
            "treatment": "IV antibiotics; spine surgery consult for abscess causing cord compression; image-guided biopsy for organism identification.",
        },
    },
    {
        "name": "Vertebral Hemangioma",
        "slides": [48],
        "modalities": ["CT", "MRI"],
        "isRedFinding": False,
        "patternName": "CT Spine T/L",
        "stepTitle": "9. Assess the osseous vertebra and discs.",
        "extra": {
            "clinical": "Incidental; polka-dot appearance on axial CT (vertical trabeculae); T1/T2 bright on MRI. Rarely aggressive.",
            "report": "Vertical trabecular pattern ('polka-dot' on axial CT) in (vertebral body) consistent with hemangioma. Appears T1/T2 bright on MRI. No aggressive features.",
            "treatment": "No further workup for typical appearance. Symptomatic/aggressive hemangiomas may require vertebroplasty or radiation.",
        },
    },
    {
        "name": "Neurogenic Tumors (Schwannoma / Neurofibroma)",
        "slides": [49],
        "modalities": ["MRI", "CT"],
        "isRedFinding": False,
        "patternName": "CT Spine T/L",
        "stepTitle": "14. Look at the paraspinal soft tissues.",
        "extra": {
            "clinical": "Posterior mediastinum or paraspinal region; dumbbell-shaped lesion extending through neural foramen.",
            "report": "Dumbbell-shaped T2-hyperintense lesion in the (right/left) paraspinal/extraforaminal region at (level) extending into the spinal canal via the neural foramen consistent with schwannoma/neurofibroma.",
            "treatment": "Spine surgery or neurosurgery consult for symptomatic lesions.",
        },
    },
    {
        "name": "Spinal Meningioma",
        "slides": [50],
        "modalities": ["MRI"],
        "isRedFinding": False,
        "patternName": "CT Spine T/L",
        "stepTitle": "5. Do a quick first look at the spinal canal.",
        "extra": {
            "clinical": "Middle-aged women; intradural extramedullary; T2 iso/hypointense with intense enhancement and dural tail.",
            "report": "Intradural extramedullary T2-isointense mass at (level) with intense homogeneous enhancement and dural tail consistent with meningioma.",
            "treatment": "Spine surgery consult; surgical resection for symptomatic lesions.",
        },
    },

    # ── HEAD AND NECK / CT SOFT TISSUE NECK ───────────────────────────────
    {
        "name": "Tonsillitis and Peritonsillar Abscess",
        "slides": [51],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CT Soft Tissue Neck",
        "stepTitle": "6. Switch to soft tissue windows to assess the mucosa and midline deep spaces.",
        "extra": {
            "clinical": "Sore throat, dysphagia, fever, muffled voice. 'Don't get worked up over tonsillar vs peritonsillar.'",
            "report": "Enlarged (tonsils/peritonsillar) with (adjacent fluid collection/rim enhancement) consistent with (tonsillitis/peritonsillar abscess). No extension to parapharyngeal or retropharyngeal space.",
            "treatment": "ENT consult; antibiotics; needle aspiration vs. incision and drainage for abscess.",
        },
    },
    {
        "name": "Masticator Space Infection",
        "slides": [52],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Soft Tissue Neck",
        "stepTitle": "7. Assess the remaining (paramedian and lateral) deep spaces of the neck.",
        "extra": {
            "clinical": "Jaw pain, trismus, facial swelling; look for odontogenic disease (cavities, periapical lucencies).",
            "report": "Soft tissue thickening and (rim-enhancing collection/phlegmon) in the masticator space consistent with masticator space abscess of likely odontogenic origin.",
            "treatment": "Oral and maxillofacial surgery/ENT consult; IV antibiotics; drainage; treat dental source.",
        },
    },
    {
        "name": "Odontogenic Subperiosteal Abscess",
        "slides": [53],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Soft Tissue Neck",
        "stepTitle": "7. Assess the remaining (paramedian and lateral) deep spaces of the neck.",
        "extra": {
            "clinical": "Dental pain, facial swelling, periapical lucency on CT. Can spread to submandibular, masticator, and parapharyngeal spaces.",
            "report": "Periapical lucency at (tooth number/location) with adjacent subperiosteal fluid collection consistent with odontogenic subperiosteal abscess.",
            "treatment": "Oral surgery/ENT consult; IV antibiotics; dental extraction and surgical drainage.",
        },
    },
    {
        "name": "Retropharyngeal Infection",
        "slides": [55],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Soft Tissue Neck",
        "stepTitle": "6. Switch to soft tissue windows to assess the mucosa and midline deep spaces.",
        "extra": {
            "clinical": "Neck pain, fever, dysphagia, stiff neck; can cause airway compromise or descend into mediastinum.",
            "report": "Fluid/rim-enhancing collection in the retropharyngeal space at (level) measuring X cm. (No extension / Extension) to the danger space. (Airway compromise present/absent.)",
            "treatment": "Airway management priority; ENT/surgery consult; IV antibiotics; surgical drainage for abscess.",
        },
    },
    {
        "name": "Laryngeal Trauma",
        "slides": [56],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Soft Tissue Neck",
        "stepTitle": "Airways , upper lungs, other structures on lung windws",
        "extra": {
            "clinical": "Direct blow to the neck; hoarseness, stridor, subcutaneous emphysema.",
            "report": "Fracture of the (thyroid/cricoid/arytenoid) cartilage with (associated soft tissue hematoma/airway compromise/subcutaneous emphysema). (Mucosal disruption noted.)",
            "treatment": "Airway management priority (secure early); ENT consult; surgical repair for displaced fractures.",
        },
    },
    {
        "name": "Esophageal Tear / Boerhaave Syndrome",
        "slides": [57],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Soft Tissue Neck",
        "stepTitle": "7. Assess the remaining (paramedian and lateral) deep spaces of the neck.",
        "extra": {
            "clinical": "Sudden vomiting followed by severe chest/back pain; Mackler triad (vomiting, lower chest pain, subcutaneous emphysema).",
            "report": "Periesophageal soft tissue gas and fluid with (mediastinal widening/pleural effusion) consistent with esophageal perforation. No contrast extravasation study performed.",
            "treatment": "Immediate surgical consult; NPO; broad-spectrum antibiotics; esophageal diversion vs. primary repair.",
        },
    },

    # ── CHEST / CT CHEST / CTA ─────────────────────────────────────────────
    {
        "name": "Pulmonary Contusion and Laceration",
        "slides": [64],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CT Chest",
        "stepTitle": "14. Examine the lungs.",
        "extra": {
            "clinical": "Blunt chest trauma; consolidation/ground-glass developing over 24-48 h. Laceration creates pneumatocele.",
            "report": "Pulmonary contusion in the (location) with adjacent laceration/pneumatocele. Estimated volume: (small/moderate/large).",
            "treatment": "Supportive; respiratory therapy; monitor for respiratory compromise.",
        },
    },
    {
        "name": "Tracheal Diverticulum",
        "slides": [65],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CT Chest",
        "stepTitle": "13. Check the airways.",
        "extra": {
            "clinical": "Incidental outpouching from the posterior tracheal wall, typically right side at T1-T2 level. Usually asymptomatic.",
            "report": "Small outpouching from the posterior wall of the trachea at the (right/left) paraspinal region consistent with a tracheal diverticulum.",
            "treatment": "No treatment needed for asymptomatic incidental finding.",
        },
    },
    {
        "name": "Pulmonary Interstitial Emphysema (PIE) and Macklin Effect",
        "slides": [66],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "13. Check the airways.",
        "extra": {
            "clinical": "Macklin effect: alveolar rupture → gas tracks along bronchovascular bundles → pneumomediastinum/PTX. Seen in both traumatic and non-traumatic settings.",
            "report": "Linear gas lucencies tracking along bronchovascular bundles consistent with PIE/Macklin effect. Associated pneumomediastinum (present/absent) and pneumothorax (present/absent).",
            "treatment": "Treat underlying cause; lung protective ventilation; may require decompression if pneumothorax develops.",
        },
    },
    {
        "name": "Mediastinal Mass (Anterior / Middle / Posterior)",
        "slides": [70, 71],
        "modalities": ["CT", "MRI"],
        "isRedFinding": False,
        "patternName": "CT Chest",
        "stepTitle": "11. Examine the Mediastinum and thoracic lymph node stations.",
        "extra": {
            "clinical": "Anterior: 4T's (Thymoma, Teratoma/germ cell, Thyroid, 'Terrible' lymphoma). Middle: vascular/cysts. Posterior: neurogenic/thoracic.",
            "report": "Soft tissue mass in the (anterior/middle/posterior) mediastinum measuring X cm. Differential: (list). MRI/PET/biopsy may be required.",
            "treatment": "Surgery/oncology/pulmonology depending on compartment and likely diagnosis.",
        },
    },
    {
        "name": "Esophageal Adenocarcinoma",
        "slides": [72],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CT Chest",
        "stepTitle": "16. Examine the esophagus from its cervical portion through the diaphragmatic",
        "extra": {
            "clinical": "Dysphagia, weight loss, GERD history; distal esophagus/GEJ. Thickening >3-5 mm is abnormal.",
            "report": "Circumferential/eccentric esophageal wall thickening at the (level/GEJ) measuring X mm. Note periesophageal fat stranding and lymphadenopathy. Staging CT recommended.",
            "treatment": "Gastroenterology/oncology consult; endoscopy/biopsy for tissue diagnosis.",
        },
    },
    {
        "name": "Esophageal Wall Thickening Mimic",
        "slides": [73],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CT Chest",
        "stepTitle": "16. Examine the esophagus from its cervical portion through the diaphragmatic",
        "extra": {
            "clinical": "Collapsed/distended esophagus can mimic wall thickening. Re-examine with distension or follow-up esophagram.",
            "report": "Apparent esophageal wall thickening may represent collapsed lumen rather than true pathology. Follow-up esophagram recommended if clinical concern persists.",
            "treatment": "Correlation with clinical context; esophagram for clarification.",
        },
    },
    {
        "name": "Esophagectomy and Gastric Pull-Through",
        "slides": [75],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CT Chest",
        "stepTitle": "16. Examine the esophagus from its cervical portion through the diaphragmatic",
        "extra": {
            "clinical": "Post-esophagectomy anatomy; gastric conduit in posterior mediastinum is normal.",
            "report": "Status post esophagectomy with gastric pull-through conduit in the posterior mediastinum. (Anastomotic leak/leak absent.) (Conduit appears patent.)",
            "treatment": "Upper GI contrast study if leak suspected; surgery consult for anastomotic complications.",
        },
    },
    {
        "name": "Benign Mediastinal Cysts (Foregut Duplication / Pericardial)",
        "slides": [76],
        "modalities": ["CT", "MRI"],
        "isRedFinding": False,
        "patternName": "CT Chest",
        "stepTitle": "11. Examine the Mediastinum and thoracic lymph node stations.",
        "extra": {
            "clinical": "Incidental thin-walled cysts; foregut duplication (posterior mediastinum); pericardial cysts (right cardiophrenic angle).",
            "report": "Thin-walled cystic structure at the (location) with water attenuation. No internal enhancement or solid components. Consistent with (foregut duplication cyst / pericardial cyst).",
            "treatment": "Observation for asymptomatic cysts; MRI for characterization if uncertain.",
        },
    },
    {
        "name": "Pulmonary Embolism (PE)",
        "slides": [77, 78],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CTA PE (CT angiogram for pulmonary embolus)",
        "stepTitle": "4. Assess the pulmonary vasculature.",
        "extra": {
            "clinical": "Acute dyspnea, plexic chest pain, hypoxia; risk factors: immobility, malignancy, hypercoagulability. Arms at side and large body habitus lower sensitivity.",
            "report": "Filling defect(s) in the (right/left) (main/lobar/segmental/subsegmental) pulmonary artery/arteries consistent with acute PE. RV:LV ratio (>1/≤1). (Pulmonary infarct present/absent.) (Pulmonary HTN features present/absent.)",
            "treatment": "Anticoagulation; consider catheter-directed thrombolysis or surgical embolectomy for massive PE with hemodynamic compromise.",
        },
    },
    {
        "name": "CTA Flow Artifact Mimicking PE",
        "slides": [79],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CTA PE (CT angiogram for pulmonary embolus)",
        "stepTitle": "4. Assess the pulmonary vasculature.",
        "extra": {
            "clinical": "Hyperdense streaks/artifacts from adjacent structures or poor bolus timing mimicking filling defects.",
            "report": "Apparent filling defect in (location) likely represents flow/beam-hardening artifact given (geometry/beam-hardening source). No definite PE identified. Repeat study with improved technique if clinical concern persists.",
            "treatment": "Repeat CTA with optimized bolus timing if clinically indicated.",
        },
    },
    {
        "name": "Septic Pulmonary Emboli",
        "slides": [80],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "14. Examine the lungs.",
        "extra": {
            "clinical": "IV drug use, right-sided endocarditis, infected IV line. Multiple peripheral nodules with cavitation.",
            "report": "Multiple bilateral peripheral cavitating pulmonary nodules consistent with septic emboli. Evaluate for right-sided cardiac vegetation on echocardiography.",
            "treatment": "IV antibiotics; cardiology consult for endocarditis evaluation; source control.",
        },
    },
    {
        "name": "Necrotizing Pneumonia and Lung Abscess",
        "slides": [81],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "14. Examine the lungs.",
        "extra": {
            "clinical": "Standard PNA in emphysema mimics necrotizing PNA. Pulmonary abscess vs. empyema: look for pulmonary parenchyma on all sides (abscess) vs. lenticular shape at pleural surface (empyema).",
            "report": "Consolidation with central low-attenuation necrosis and (cavitation/air-fluid level) consistent with necrotizing pneumonia / lung abscess.",
            "treatment": "IV antibiotics; CT-guided or bronchoscopic drainage for abscess; surgery for refractory cases.",
        },
    },
    {
        "name": "Aortic Injury / Traumatic Aortic Injury",
        "slides": [82, 83, 84],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "10. Assess the vasculature.",
        "extra": {
            "clinical": "High-speed deceleration; most at aortic isthmus. Stanford A: proximal to left subclavian. Stanford B: no proximal aorta involvement. New SVS/STS classification uses innominate artery as A/B border.",
            "report": "Mediastinal hematoma with (intimal irregularity / pseudoaneurysm / transection) at the aortic isthmus consistent with traumatic aortic injury. SVS Grade (I–IV).",
            "treatment": "Vascular surgery/interventional radiology consult; TEVAR for B injuries; open repair for A injuries; permissive hypotension until repair.",
        },
    },
    {
        "name": "Aortic Dissection (Stanford Classification)",
        "slides": [83, 84],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CTA Aorta (Dissection Study)",
        "stepTitle": "5. Examine the contrast enhanced study, starting with the aorta.",
        "extra": {
            "clinical": "Sudden tearing chest/back pain; hypertension. Stanford A: involves ascending aorta (surgical emergency). Stanford B: no ascending involvement (medical management or TEVAR).",
            "report": "Type (A/B) aortic dissection with intimal flap from (origin) to (extent). True lumen: (compressed/patent). Involvement of (vessels). (Branch vessel compromise identified.)",
            "treatment": "Type A: emergent cardiac surgery. Type B: IV antihypertensive; TEVAR for complicated B.",
        },
    },
    {
        "name": "Traumatic Pneumothorax",
        "slides": [85],
        "modalities": ["CT", "Plain Radiograph"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "14. Examine the pleura.",
        "extra": {
            "clinical": "Chest trauma; deep sulcus sign on supine CXR. Identify the visceral pleural line on CT.",
            "report": "Pneumothorax along the (anterior/lateral/apical) pleural surface measuring approximately X cm. (Tension features absent/present: mediastinal shift, diaphragm inversion.)",
            "treatment": "Chest tube for moderate-large PTX or tension; observation for small asymptomatic PTX.",
        },
    },
    {
        "name": "Spontaneous Pneumothorax",
        "slides": [86],
        "modalities": ["CT", "Plain Radiograph"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "14. Examine the pleura.",
        "extra": {
            "clinical": "No trauma history; look for underlying bleb/bullae. Primary (young/thin) vs. secondary (COPD, Marfan).",
            "report": "Spontaneous pneumothorax with (bleb/bulla) at the (lung apex). No traumatic etiology identified.",
            "treatment": "Chest tube; treat underlying cause; thoracic surgery consult for recurrent/bilateral.",
        },
    },
    {
        "name": "Flail Chest",
        "slides": [87],
        "modalities": ["CT", "Plain Radiograph"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "8. Assess the osseous structures.",
        "extra": {
            "clinical": "Three or more consecutive ribs fractured in two places; paradoxical movement. Often associated with pulmonary contusion.",
            "report": "Multiple (consecutive) rib fractures at (levels/locations) with (anterior/posterior/lateral) flail segment. Associated pulmonary contusion (present/absent).",
            "treatment": "Aggressive pain management; positive pressure ventilation; surgical rib fixation for refractory respiratory failure.",
        },
    },
    {
        "name": "Sternal Fractures",
        "slides": [88],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "8. Assess the osseous structures.",
        "extra": {
            "clinical": "Commonly missed. Look at the sternum on sagittal view on every trauma scan, especially if no clear anterior chest injury history.",
            "report": "Fracture of the (manubrium/body) of the sternum with (displacement/no displacement). (Retrosternal hematoma present/absent.)",
            "treatment": "Cardiac monitoring (sternal fracture associated with myocardial contusion); surgical fixation for significant displacement.",
        },
    },
    {
        "name": "Diaphragm Rupture",
        "slides": [89],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "14. Examine the pleura.",
        "extra": {
            "clinical": "Blunt abdominal trauma; 'fallen viscus' or 'dependent viscera' sign — bowel/stomach herniated into chest and falls dependently.",
            "report": "Herniation of (stomach/bowel) through the (left/right) hemidiaphragm with diaphragmatic discontinuity consistent with diaphragmatic rupture.",
            "treatment": "Surgical repair; nasogastric tube decompression.",
        },
    },
    {
        "name": "Pleural Plaques (Asbestos-Related Pleural Disease)",
        "slides": [90],
        "modalities": ["CT", "Plain Radiograph"],
        "isRedFinding": False,
        "patternName": "CT Chest",
        "stepTitle": "14. Examine the pleura.",
        "extra": {
            "clinical": "Discontinuous calcified pleural plaques = asbestos-related pleural disease, NOT asbestosis (which is the interstitial lung disease).",
            "report": "Bilateral discontinuous calcified pleural plaques consistent with asbestos-related pleural disease. No evidence of malignant pleural thickening.",
            "treatment": "Surveillance; mesothelioma screening in high-risk patients.",
        },
    },
    {
        "name": "Pleural Metastases",
        "slides": [91],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "14. Examine the pleura.",
        "extra": {
            "clinical": "Nodular pleural thickening, malignant effusion; common primaries: lung, breast, ovarian, mesothelioma.",
            "report": "Nodular/irregular pleural thickening with (hemorrhagic/exudative) pleural effusion consistent with pleural metastases.",
            "treatment": "Oncology consult; thoracentesis for cytology; pleurodesis for recurrent effusion.",
        },
    },
    {
        "name": "Empyema Necessitatis",
        "slides": [92],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "14. Examine the pleura.",
        "extra": {
            "clinical": "Infected pleural fluid erodes through pleura and parietal chest wall into soft tissues.",
            "report": "Loculated pleural collection with extension through the chest wall into the soft tissues consistent with empyema necessitatis.",
            "treatment": "Cardiothoracic surgery consult; drainage; IV antibiotics.",
        },
    },
    {
        "name": "Elastofibroma Dorsi",
        "slides": [93],
        "modalities": ["CT", "MRI"],
        "isRedFinding": False,
        "patternName": "CT Chest",
        "stepTitle": "7. Evaluate the musculature.",
        "extra": {
            "clinical": "Elderly; bilateral subscapular soft tissue masses with fat strands. Don't call it a hematoma and don't recommend further imaging.",
            "report": "Bilateral subscapular soft tissue masses with interspersed fat strands consistent with elastofibroma dorsi.",
            "treatment": "No further workup needed for classic appearance.",
        },
    },

    # ── PEDIATRIC ─────────────────────────────────────────────────────────
    {
        "name": "Foreign Body Inhalation",
        "slides": [120],
        "modalities": ["Plain Radiograph", "CT"],
        "isRedFinding": True,
        "patternName": "CT Chest",
        "stepTitle": "13. Check the airways.",
        "extra": {
            "clinical": "Young child; sudden coughing/choking; asymmetric hyperinflation on expiratory CXR.",
            "report": "Radiopaque/radiolucent foreign body in the (trachea/right/left) mainstem bronchus / (location). Air trapping in the (right/left) lung field consistent with partial obstruction.",
            "treatment": "Pediatric pulmonology/surgery consult; rigid bronchoscopy for retrieval.",
        },
    },
    {
        "name": "Intussusception",
        "slides": [121],
        "modalities": ["CT", "US"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Ileocolic (>2.5 cm wall-to-wall, >3.5 cm length, RLQ) — call peds attending for reduction. SB-SB (<2.5 cm) often transient/benign. Target sign on US.",
            "report": "Ileocolic intussusception with (lead point identified / no lead point). Bowel wall-to-wall diameter X cm, length X cm. (Ischemia features absent/present.)",
            "treatment": "Ileocolic: enema reduction (air/water) under fluoroscopy or US. SB-SB: observation or surgical evaluation.",
        },
    },

    # ── BODY / CT ABDOMEN PELVIS ──────────────────────────────────────────
    {
        "name": "Hypoperfusion Complex (Shock Bowel)",
        "slides": [123],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "8. Assess the vasculature.",
        "extra": {
            "clinical": "Severe hypotension/shock; diffuse small bowel hyperperfusion, flat IVC, small aorta, adrenal/pancreatic enhancement.",
            "report": "Diffuse small bowel wall thickening and hyperperfusion with flat IVC and small aortic caliber consistent with hypoperfusion complex. Correlate with hemodynamic status.",
            "treatment": "Aggressive resuscitation; trauma surgery consultation.",
        },
    },
    {
        "name": "Bowel Injury",
        "slides": [124],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Blunt abdominal trauma; free air/fluid without solid organ injury suggests bowel injury.",
            "report": "Free intraperitoneal air adjacent to the (small/large) bowel with associated mesenteric hematoma and bowel wall thickening consistent with bowel perforation.",
            "treatment": "Emergent surgery consult.",
        },
    },
    {
        "name": "Intraperitoneal Infection / Peritonitis",
        "slides": [125],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Post-operative fluid collections in surgical beds. Free air without surgical history = perforation until proven otherwise.",
            "report": "Complex fluid collection at the (surgical bed/location) with rim enhancement and adjacent fat stranding consistent with postoperative abscess / peritonitis.",
            "treatment": "Surgery/IR consult; image-guided drainage or surgical exploration.",
        },
    },
    {
        "name": "External Hernias (Inguinal / Femoral / Obturator)",
        "slides": [126],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Groin bulge; femoral hernia in elderly females (below inguinal ligament). Obturator hernia: elderly women, Howship-Romberg sign.",
            "report": "Inguinal/femoral/obturator hernia containing (bowel/omentum) with (reducible/incarcerated/strangulated) appearance. (Signs of bowel obstruction/ischemia: present/absent.)",
            "treatment": "Surgical repair urgently for strangulated hernias.",
        },
    },
    {
        "name": "Hiatus Hernia",
        "slides": [127],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Sliding (type I) vs. paraesophageal (type II–IV). Large hernias can volvulize.",
            "report": "Hiatus hernia type (I/II/III/IV) with herniation of (fundus/entire stomach/other organs) into the thoracic cavity.",
            "treatment": "Observe type I; surgery for symptomatic paraesophageal hernias.",
        },
    },
    {
        "name": "Gastritis and Perforated Peptic Ulcer",
        "slides": [128],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Epigastric pain; perforation presents with peritonitis. Caustic ingestion causes similar diffuse gastric wall thickening.",
            "report": "Gastric wall thickening with pneumoperitoneum and perihepatic/periesophageal free air consistent with perforated peptic ulcer.",
            "treatment": "Surgery consult for perforation; PPI therapy for gastritis.",
        },
    },
    {
        "name": "Gastric Malignancy",
        "slides": [129],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Dysphagia, weight loss, early satiety. Circumferential wall thickening, lymphadenopathy.",
            "report": "Circumferential gastric wall thickening at the (antrum/fundus/body) with perigastric lymphadenopathy consistent with gastric carcinoma. Staging CT recommended.",
            "treatment": "Gastroenterology/oncology consult; endoscopy and biopsy.",
        },
    },
    {
        "name": "Gastric Volvulus",
        "slides": [130],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Organoaxial (more common) vs. mesenteroaxial. Can cause gastric ischemia.",
            "report": "Gastric volvulus (organoaxial/mesenteroaxial) with abnormal position of the stomach. (Ischemia features: present/absent.)",
            "treatment": "Emergent surgery for volvulus with ischemia; endoscopic decompression vs. surgery.",
        },
    },
    {
        "name": "Duodenal Hematoma",
        "slides": [131],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Blunt abdominal trauma to the right upper quadrant; handlebar injury in children. Associated with pancreatic injury.",
            "report": "Intramural duodenal hematoma causing partial obstruction. Assess for associated pancreatic injury.",
            "treatment": "Conservative management with NGT decompression; surgery for complete obstruction or perforation.",
        },
    },
    {
        "name": "Bowel Ischemia",
        "slides": [137],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Bowel wall thickening, pneumatosis intestinalis, portal venous gas, hypo-enhancing bowel. Reperfusion after ischemia can appear hyperperfused.",
            "report": "Bowel wall thickening with (pneumatosis / portal venous gas / mesenteric stranding) consistent with bowel ischemia. (Free perforation: present/absent.)",
            "treatment": "Emergent surgery consult; broad spectrum antibiotics.",
        },
    },
    {
        "name": "Infectious Colitis",
        "slides": [138],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Diarrhea, fever; C. diff most common in hospitalized patients.",
            "report": "Diffuse colonic wall thickening with (pericolonic fat stranding/ascites/toxic dilation) consistent with (infectious/C. diff) colitis.",
            "treatment": "IV fluids, antibiotics; surgery for toxic megacolon or perforation.",
        },
    },
    {
        "name": "Typhlitis / Neutropenic Colitis",
        "slides": [139],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Neutropenic patient (chemotherapy); right hemicolon inflammation. High mortality.",
            "report": "Right hemicolonic wall thickening with pericolonic fat stranding in a neutropenic patient consistent with typhlitis/neutropenic colitis.",
            "treatment": "Bowel rest, IV antibiotics, G-CSF; surgery for perforation.",
        },
    },
    {
        "name": "Ulcerative Colitis / Toxic Megacolon",
        "slides": [140],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Toxic megacolon: dilated colon with air-fluid levels, thin wall, pseudopolyps; triggered by C. diff or UC flare.",
            "report": "Toxic megacolon with colonic dilation (>6 cm transverse) and (thin wall / mucosal pseudopolyps / thumb-printing). High risk of perforation.",
            "treatment": "Bowel rest, IV steroids, IV antibiotics; colectomy for perforation or clinical deterioration.",
        },
    },
    {
        "name": "Appendicitis",
        "slides": [141, 142],
        "modalities": ["CT", "US", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "RLQ pain, fever, leukocytosis. Dilated appendix >6 mm with periappendiceal fat stranding. Appendicolith → higher risk of perforation. DWI can help on MRI.",
            "report": "Dilated appendix (X mm) with periappendiceal fat stranding consistent with acute appendicitis. (Appendicolith present/absent.) (Perforation/abscess: present/absent.)",
            "treatment": "Surgery consult; laparoscopic appendectomy; IV antibiotics for complicated cases.",
        },
    },
    {
        "name": "Epiploic Appendagitis and Omental Infarction",
        "slides": [143],
        "modalities": ["CT"],
        "isRedFinding": False,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Epiploic appendagitis (EA): focal LLQ (classically). Self-limiting. Torsed appendages calcify chronically (explain weird benign calcifications).",
            "report": "Fatty lesion with surrounding fat stranding and central dense dot adjacent to the (sigmoid) colon consistent with epiploic appendagitis. No bowel obstruction.",
            "treatment": "NSAIDs; self-limiting; no surgery needed.",
        },
    },
    {
        "name": "Colorectal Carcinoma",
        "slides": [145],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "10. Assess the Gastrointestinal tract.",
        "extra": {
            "clinical": "Change in bowel habits, rectal bleeding; 'apple-core' lesion or circumferential wall thickening.",
            "report": "Focal circumferential wall thickening of the (location) with associated (lymphadenopathy/liver lesions) consistent with colorectal carcinoma. Staging CT recommended.",
            "treatment": "Colorectal surgery/oncology consult; colonoscopy and biopsy.",
        },
    },
    {
        "name": "Splenic Abscess",
        "slides": [147],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "14. Examine the spleen.",
        "extra": {
            "clinical": "Fever, LUQ pain; immunocompromised, endocarditis, or trauma. Pyogenic vs. microabscesses (fungal/TB).",
            "report": "Low-attenuation (rim-enhancing) lesion(s) in the spleen consistent with splenic abscess.",
            "treatment": "IV antibiotics; percutaneous drainage or splenectomy depending on size and response.",
        },
    },
    {
        "name": "Splenic Infarct",
        "slides": [148],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "14. Examine the spleen.",
        "extra": {
            "clinical": "LUQ pain; causes: embolic (atrial fibrillation, endocarditis), hypercoagulable, sickle cell. Ensure this isn't a normal finding from early contrast bolus (wedge-shaped defects in arterial phase can be normal).",
            "report": "Wedge-shaped non-enhancing splenic parenchymal defect consistent with splenic infarct.",
            "treatment": "Treat underlying cause; anticoagulation if embolic; splenectomy for hemorrhagic infarct.",
        },
    },
    {
        "name": "Liver Abscess",
        "slides": [150],
        "modalities": ["CT", "US"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "11. Assess the liver.",
        "extra": {
            "clinical": "RUQ pain, fever, elevated WBC; pyogenic (most common) vs. amoebic.",
            "report": "Hypodense rim-enhancing hepatic collection in (segment) consistent with hepatic abscess.",
            "treatment": "IV antibiotics; percutaneous drainage; surgery for multiloculated or non-drainable abscesses.",
        },
    },
    {
        "name": "Hepatitis (Acute)",
        "slides": [151],
        "modalities": ["CT", "US"],
        "isRedFinding": False,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "11. Assess the liver.",
        "extra": {
            "clinical": "Elevated LFTs, RUQ pain/tenderness; periportal oedema on CT/US.",
            "report": "Hepatomegaly with periportal oedema and gallbladder wall thickening consistent with acute hepatitis.",
            "treatment": "Hepatology consult; supportive care; treat underlying cause.",
        },
    },
    {
        "name": "Hepatic Steatosis (Fatty Liver)",
        "slides": [152],
        "modalities": ["CT", "US", "MRI"],
        "isRedFinding": False,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "11. Assess the liver.",
        "extra": {
            "clinical": "Incidental on CT; liver attenuation <40 HU or ≥10 HU less than spleen. Geographic sparing common around gallbladder fossa and ligamentum teres.",
            "report": "Hepatic steatosis with liver attenuation less than spleen. Geographic areas of fatty sparing noted. No focal liver lesion.",
            "treatment": "Lifestyle modification; hepatology referral for significant steatosis.",
        },
    },
    {
        "name": "Ascending Cholangitis",
        "slides": [153],
        "modalities": ["CT", "US"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "12. Assess the gallbladder and remaining biliary tree.",
        "extra": {
            "clinical": "Charcot's triad (fever, RUQ pain, jaundice); Reynolds' pentad adds hypotension and altered mental status.",
            "report": "Biliary wall thickening/hyperenhancement with biliary debris and dilated bile ducts consistent with ascending cholangitis. Choledocholithiasis (present/absent).",
            "treatment": "IV antibiotics; urgent ERCP for biliary decompression; surgery for non-responsive cases.",
        },
    },
    {
        "name": "Pancreatitis (Acute)",
        "slides": [156],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "13. Examine the pancreas.",
        "extra": {
            "clinical": "Epigastric pain radiating to back; elevated lipase. Gallstone or alcohol most common causes.",
            "report": "Pancreatic enlargement with peripancreatic fat stranding (CT Severity Index). (Necrosis: absent/present — affecting X% of gland.) (Fluid collections: absent/present.)",
            "treatment": "IV fluids, bowel rest, pain management; treat underlying cause; ICU for severe pancreatitis.",
        },
    },
    {
        "name": "Pancreatic Necrosis and Walled-Off Necrosis",
        "slides": [157],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "13. Examine the pancreas.",
        "extra": {
            "clinical": "Complication of pancreatitis; peripancreatic necrosis evolves to walled-off necrosis (WON) after 4 weeks.",
            "report": "Pancreatic/peripancreatic necrosis involving (X%) of the gland with (gas) indicating infection. Walled-off necrosis collection at (location) measuring X cm.",
            "treatment": "Interventional gastroenterology/surgery consult; antibiotics; endoscopic or percutaneous drainage of infected WON.",
        },
    },
    {
        "name": "Pancreatic Laceration",
        "slides": [159],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "13. Examine the pancreas.",
        "extra": {
            "clinical": "Blunt upper abdominal trauma; handlebar injury. Pancreatic duct injury determines management.",
            "report": "Pancreatic laceration at the (head/neck/body/tail) with (associated peripancreatic fluid). Main pancreatic duct (intact/disrupted) — MRCP recommended for duct evaluation.",
            "treatment": "Surgery consult; MRCP or ERCP for duct injury assessment; distal pancreatectomy for distal duct disruption.",
        },
    },
    {
        "name": "Pancreatic Adenocarcinoma",
        "slides": [160],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "13. Examine the pancreas.",
        "extra": {
            "clinical": "Weight loss, painless jaundice; 'double duct sign'. Most at pancreatic head.",
            "report": "Hypoenhancing mass at the (head/body/tail) of the pancreas causing (main pancreatic duct / common bile duct) dilation. (Vascular encasement/contact: present/absent.) Staging suggests (resectable/borderline/unresectable).",
            "treatment": "Hepatobiliary surgery/oncology consult; Whipple procedure for resectable head tumors.",
        },
    },
    {
        "name": "Retroperitoneal Hemorrhage",
        "slides": [161],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "20. Assess the lymphatic system, mesentery, and other potential spaces.",
        "extra": {
            "clinical": "Spontaneous (anticoagulation, warfarin) or traumatic. Often with layering hematocrit level.",
            "report": "Retroperitoneal hemorrhage with hematocrit layering. (Active extravasation: present/absent.) (Psoas hematoma extending to (location).)",
            "treatment": "Reverse anticoagulation; IR for embolization of active bleeding; surgery for unstable patients.",
        },
    },
    {
        "name": "Retroperitoneal Lymphoma and Metastases",
        "slides": [162],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "20. Assess the lymphatic system, mesentery, and other potential spaces.",
        "extra": {
            "clinical": "Lymphoma: 'rubber-hose' sign; encases vessels without invasion. Metastatic lymph nodes common.",
            "report": "Retroperitoneal lymphadenopathy with nodal masses measuring up to X cm consistent with (lymphoma/metastatic disease). (Vascular encasement without invasion.)",
            "treatment": "Oncology/hematology consult; PET-CT for staging; biopsy.",
        },
    },
    {
        "name": "Adrenal Hemorrhage",
        "slides": [164],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "15. Examine the adrenal glands.",
        "extra": {
            "clinical": "Trauma, neonates, sepsis, anticoagulation. Bilateral adrenal hemorrhage → adrenal insufficiency.",
            "report": "Hyperdense (pre-contrast) or heterogeneous adrenal gland enlargement at the (right/left/bilateral) consistent with adrenal hemorrhage.",
            "treatment": "Endocrinology consult for bilateral hemorrhage; stress-dose steroids if adrenal insufficiency present.",
        },
    },
    {
        "name": "Adrenal Adenoma",
        "slides": [165],
        "modalities": ["CT", "MRI"],
        "isRedFinding": False,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "15. Examine the adrenal glands.",
        "extra": {
            "clinical": "Most common adrenal mass; lipid-rich (<10 HU pre-contrast) or lipid-poor but washout >40% on delayed imaging.",
            "report": "Adrenal nodule measuring X mm with attenuation of X HU (pre-contrast) consistent with lipid-rich adenoma. No malignant features.",
            "treatment": "Follow ACR Incidentaloma guidelines; functional workup if clinically indicated.",
        },
    },
    {
        "name": "Pyelonephritis",
        "slides": [167],
        "modalities": ["CT", "US", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "16. Examine the kidneys.",
        "extra": {
            "clinical": "Flank pain, fever, dysuria; DWI and ADC most sensitive on MRI.",
            "report": "Wedge-shaped hypoenhancing zone(s) in the (right/left) kidney with perinephric fat stranding consistent with acute pyelonephritis. (Abscess: present/absent.)",
            "treatment": "IV antibiotics; urology consult for obstruction or abscess.",
        },
    },
    {
        "name": "Emphysematous Pyelonephritis / Emphysematous Pyelitis",
        "slides": [168],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "16. Examine the kidneys.",
        "extra": {
            "clinical": "Emphysematous pyelitis: gas in collecting system (less severe). Emphysematous pyelonephritis: gas in renal parenchyma (very severe, high mortality). Diabetics at highest risk.",
            "report": "Gas within the (collecting system: pyelitis / renal parenchyma: pyelonephritis) of the (right/left) kidney consistent with emphysematous pyelitis/pyelonephritis.",
            "treatment": "Emphysematous pyelonephritis: urgent urology consult; IV antibiotics; percutaneous drainage vs. nephrectomy.",
        },
    },
    {
        "name": "Bladder Rupture (Intraperitoneal and Extraperitoneal)",
        "slides": [172],
        "modalities": ["CT"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "18. Check the bladder.",
        "extra": {
            "clinical": "Pelvic trauma with pelvic fracture. Extraperitoneal: more common, associated with pelvic fracture. Intraperitoneal: from direct blow to distended bladder.",
            "report": "CT cystogram demonstrates (intraperitoneal/extraperitoneal) contrast leak from the bladder consistent with bladder rupture.",
            "treatment": "Extraperitoneal: Foley catheter drainage. Intraperitoneal: surgical repair.",
        },
    },
    {
        "name": "Renal Cell Carcinoma and Transitional Cell Carcinoma",
        "slides": [173],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "16. Examine the kidneys.",
        "extra": {
            "clinical": "RCC: hematuria, flank pain, palpable mass; clear cell most common. TCC: urothelial thickening, filling defect.",
            "report": "Enhancing renal mass at (location) measuring X cm consistent with RCC. (Perinephric extension/vascular involvement/lymphadenopathy: present/absent.) Stage: (T1a/T1b/T2/T3/T4).",
            "treatment": "Urology/oncology consult; partial vs. radical nephrectomy; ablation for small tumors.",
        },
    },
    {
        "name": "Bladder Cancer",
        "slides": [174],
        "modalities": ["CT", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Abdomen Pelvis",
        "stepTitle": "18. Check the bladder.",
        "extra": {
            "clinical": "Painless hematuria; most are TCC. CT urogram for staging.",
            "report": "Intraluminal bladder mass with focal wall thickening at the (posterior wall/trigone/dome) consistent with bladder carcinoma. (Perivesical extension/lymphadenopathy: present/absent.)",
            "treatment": "Urology consult; cystoscopy and biopsy; TURBT vs. cystectomy based on stage.",
        },
    },
    {
        "name": "Perigestational / Subchorionic Hematoma",
        "slides": [183],
        "modalities": ["US", "MRI"],
        "isRedFinding": True,
        "patternName": "CT Pelvis",
        "stepTitle": "13. Assess the remaining pelvic contents, i.e. genitourinary tract as well as the",
        "extra": {
            "clinical": "Vaginal bleeding in first trimester; crescent-shaped collection between chorion and uterine wall.",
            "report": "Subchorionic hematoma (small/moderate/large) measuring X cm with (acute/subacute) blood products. Fetal cardiac activity (present/absent).",
            "treatment": "Obstetrics consult; pelvic rest; close follow-up.",
        },
    },
    {
        "name": "Ectopic Pregnancy",
        "slides": [187],
        "modalities": ["US"],
        "isRedFinding": True,
        "patternName": "CT Pelvis",
        "stepTitle": "13. Assess the remaining pelvic contents, i.e. genitourinary tract as well as the",
        "extra": {
            "clinical": "Amenorrhea + pelvic pain + vaginal bleeding + positive β-hCG. Empty uterus with adnexal mass → ectopic until proven otherwise.",
            "report": "Empty uterine cavity with (right/left) adnexal gestational sac/ring/mass. Free fluid in the pelvis consistent with ruptured ectopic pregnancy.",
            "treatment": "Emergent OB/GYN consult; methotrexate vs. surgical management depending on stability and findings.",
        },
    },
    {
        "name": "Endometrial Carcinoma / Cervical Carcinoma / Ovarian Carcinoma",
        "slides": [194],
        "modalities": ["CT", "MRI", "US"],
        "isRedFinding": True,
        "patternName": "CT Pelvis",
        "stepTitle": "13. Assess the remaining pelvic contents, i.e. genitourinary tract as well as the",
        "extra": {
            "clinical": "Postmenopausal bleeding (endometrial); cervical mass on exam; ovarian pelvic mass with ascites.",
            "report": "Uterine/cervical/ovarian mass with (pelvic lymphadenopathy/peritoneal implants/ascites) consistent with gynecologic malignancy. MRI pelvis recommended for staging.",
            "treatment": "OB/GYN oncology consult; MRI staging; surgical vs. chemoradiation based on type and stage.",
        },
    },
    {
        "name": "Hemorrhagic Ovarian Cyst",
        "slides": [196],
        "modalities": ["US", "CT", "MRI"],
        "isRedFinding": False,
        "patternName": "CT Pelvis",
        "stepTitle": "13. Assess the remaining pelvic contents, i.e. genitourinary tract as well as the",
        "extra": {
            "clinical": "Premenopausal; pelvic pain; can cause hemoperitoneum if ruptured. Lace-like fibrin strands on US.",
            "report": "Complex adnexal cyst with internal echoes/blood products consistent with hemorrhagic ovarian cyst. (No solid component identified.) Follow-up US in 6 weeks recommended.",
            "treatment": "Supportive; OB/GYN consult if rupture with significant hemoperitoneum.",
        },
    },
    {
        "name": "PID and Tubo-Ovarian Abscess",
        "slides": [198],
        "modalities": ["CT", "US"],
        "isRedFinding": True,
        "patternName": "CT Pelvis",
        "stepTitle": "13. Assess the remaining pelvic contents, i.e. genitourinary tract as well as the",
        "extra": {
            "clinical": "Pelvic pain, vaginal discharge, fever in sexually active woman; cervical motion tenderness.",
            "report": "Complex cystic adnexal mass with rim enhancement consistent with tubo-ovarian abscess. Pyosalpinx (present/absent). Free pelvic fluid.",
            "treatment": "OB/GYN consult; IV antibiotics; drainage for large abscess.",
        },
    },

    # ── BONE / MSK ────────────────────────────────────────────────────────
    {
        "name": "Pathologic Fracture",
        "slides": [200],
        "modalities": ["CT", "Plain Radiograph", "MRI"],
        "isRedFinding": True,
        "patternName": "Any Bone Radiographs",
        "stepTitle": "4. Look for fractures.",
        "extra": {
            "clinical": "Fracture through abnormal bone; lytic or sclerotic lesion at fracture site; often minimal trauma.",
            "report": "Fracture through a lytic/sclerotic lesion at (location) consistent with a pathologic fracture. Underlying osseous lesion requires further workup with (MRI/bone scan/CT).",
            "treatment": "Orthopedics/oncology consult; MRI for marrow characterization; staging.",
        },
    },
    {
        "name": "Osteomyelitis",
        "slides": [201],
        "modalities": ["MRI", "CT", "Plain Radiograph"],
        "isRedFinding": True,
        "patternName": "Any Bone Radiographs",
        "stepTitle": "4. Look for fractures.",
        "extra": {
            "clinical": "Bone pain, fever, elevated inflammatory markers; hematogenous or contiguous spread.",
            "report": "Marrow signal abnormality in (bone/location) with cortical destruction and periosteal reaction consistent with osteomyelitis. Soft tissue extension (present/absent).",
            "treatment": "IV antibiotics; orthopedics consult; surgical debridement for chronic or refractory cases.",
        },
    },
    {
        "name": "Neuropathic Arthropathy (Charcot Joint)",
        "slides": [202],
        "modalities": ["CT", "Plain Radiograph", "MRI"],
        "isRedFinding": False,
        "patternName": "Any Bone Radiographs",
        "stepTitle": "9. Overall bone quality and morphology",
        "extra": {
            "clinical": "Diabetes, syphilis, syringomyelia; painless joint destruction. The 5 D's: Density, Debris, Destruction, Disorganization, Dislocation.",
            "report": "Destructive arthropathy of the (ankle/tarsal/other) joint with osseous debris, soft tissue swelling, and disorganization consistent with neuropathic/Charcot arthropathy.",
            "treatment": "Orthopedics consult; total contact cast; surgical stabilization for instability.",
        },
    },
]


def extract_slide_images(slide, max_width=800, jpeg_quality=75):
    """Return a list of base64-encoded JPEG images from the slide."""
    images = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            img_bytes = shape.image.blob
            img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
            # Resize if too large to keep JSON size manageable
            if img.width > max_width:
                ratio = max_width / img.width
                new_size = (max_width, int(img.height * ratio))
                img = img.resize(new_size, Image.Resampling.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=jpeg_quality, optimize=True)
            b64 = base64.b64encode(buf.getvalue()).decode("ascii")
            images.append(b64)
        except Exception as exc:
            print(f"  [warn] could not extract image from slide: {exc}", file=sys.stderr)
    return images


def extract_slide_text(slide):
    """Return all non-empty text lines from a slide."""
    lines = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:
                lines.append(text)
    return lines


def build_rich_content(text_lines, images):
    """Build the rich-content array for a finding."""
    content = []
    if text_lines:
        body = "\n".join(text_lines)
        content.append({"type": "text", "text": body, "bold": False, "color": None})
    for b64 in images:
        content.append({"type": "image", "format": "jpeg", "data": b64})
    return content


def main():
    pptx_path = os.path.abspath(PPTX_PATH)
    if not os.path.exists(pptx_path):
        sys.exit(f"ERROR: PPTX not found at {pptx_path}")

    print(f"Opening: {pptx_path}")
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    print(f"  {len(slides)} slides found")

    # Pre-extract all slide content
    slide_texts = {}
    slide_images = {}
    for idx, slide in enumerate(slides):
        slide_num = idx + 1
        slide_texts[slide_num] = extract_slide_text(slide)
        slide_images[slide_num] = extract_slide_images(slide)
        total_imgs = len(slide_images[slide_num])
        if total_imgs:
            print(f"  Slide {slide_num}: {total_imgs} image(s) extracted")

    # Build output seed
    seed = []
    for spec in FINDINGS_SPEC:
        all_text = []
        all_images = []
        for sn in spec["slides"]:
            if sn > len(slides):
                print(f"  [warn] slide {sn} out of range for finding '{spec['name']}'")
                continue
            all_text.extend(slide_texts[sn])
            all_images.extend(slide_images[sn])

        # Deduplicate text lines preserving order
        seen_text = set()
        unique_text = []
        for line in all_text:
            if line not in seen_text:
                seen_text.add(line)
                unique_text.append(line)

        extra = spec.get("extra", {})
        rich_content = build_rich_content(unique_text, all_images)

        entry = {
            "name": spec["name"],
            "modalities": spec["modalities"],
            "isRedFinding": spec["isRedFinding"],
            "patternName": spec["patternName"],
            "stepTitle": spec["stepTitle"],
            "clinical": extra.get("clinical", ""),
            "imaging": " ".join(unique_text) if unique_text else "",
            "report": extra.get("report", ""),
            "treatment": extra.get("treatment", ""),
            "content": rich_content,
        }
        seed.append(entry)
        print(f"  [{spec['patternName']}] {spec['name']}: "
              f"{len(unique_text)} text lines, {len(all_images)} image(s)")

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    with open(OUT_PATH, "w", encoding="utf-8") as fh:
        json.dump(seed, fh, ensure_ascii=False, indent=2)

    print(f"\nWrote {len(seed)} findings → {OUT_PATH}")
    total_with_images = sum(1 for e in seed if any(c["type"] == "image" for c in e["content"]))
    print(f"  {total_with_images} findings include at least one image")


if __name__ == "__main__":
    main()
